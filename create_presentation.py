# create_presentation.py
"""
Run this script to generate the PowerPoint presentation:
    pip install python-pptx
    python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ══════════════════════════════════════════════════════
# COLOUR PALETTE
# ══════════════════════════════════════════════════════
C_BG         = RGBColor(0x07, 0x0D, 0x1A)   # deep navy black
C_CARD       = RGBColor(0x0F, 0x1D, 0x35)   # dark navy card
C_CARD2      = RGBColor(0x13, 0x25, 0x44)   # slightly lighter card
C_ACCENT     = RGBColor(0x00, 0xC2, 0xFF)   # electric cyan
C_ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)   # vivid purple
C_GREEN      = RGBColor(0x00, 0xE0, 0x7A)   # neon green
C_ORANGE     = RGBColor(0xFF, 0xA0, 0x30)   # warm orange
C_RED        = RGBColor(0xFF, 0x4C, 0x4C)   # alert red
C_YELLOW     = RGBColor(0xFF, 0xD6, 0x00)   # highlight yellow
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT      = RGBColor(0xC8, 0xD8, 0xF0)   # soft blue-white
C_MID        = RGBColor(0x55, 0x6B, 0x8A)   # mid grey-blue
C_DARK_STRIP = RGBColor(0x03, 0x07, 0x10)   # footer strip

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

TOTAL_SLIDES = 16


# ══════════════════════════════════════════════════════
# PRIMITIVE HELPERS
# ══════════════════════════════════════════════════════

def rect(sl, l, t, w, h, rgb, line=False, line_rgb=None, line_w=Pt(1)):
    s = sl.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    if line and line_rgb:
        s.line.color.rgb = line_rgb
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def txb(sl, text, l, t, w, h,
        size=16, bold=False, italic=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = sl.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_line(sl, items, tf, size=14,
             color=C_LIGHT, gap=4, bold=False):
    """Append bullet lines to an existing text-frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(gap)
        p.alignment    = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = item
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = color


def bullet_box(sl, items, l, t, w, h,
               size=14, color=C_LIGHT, gap=4):
    box = sl.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    add_line(sl, items, tf, size=size, color=color, gap=gap)
    return box


def slide_bg(sl):
    rect(sl, 0, 0, 13.33, 7.5, C_BG)


def h_line(sl, y, color=C_ACCENT, h=0.04):
    rect(sl, 0, y, 13.33, h, color)


def slide_num(sl, n):
    txb(sl, f"{n} / {TOTAL_SLIDES}",
        11.6, 7.1, 1.6, 0.32,
        size=10, color=C_MID, align=PP_ALIGN.RIGHT)


def chip(sl, label, l, t,
         bg_col=C_ACCENT2, fg_col=C_WHITE, size=10):
    cw = len(label) * 0.10 + 0.25
    rect(sl, l, t, cw, 0.27, bg_col)
    txb(sl, label, l+0.07, t+0.03, cw-0.1, 0.22,
        size=size, bold=True, color=fg_col,
        align=PP_ALIGN.CENTER)
    return cw


def section_tag(sl, label, col=C_ACCENT2):
    chip(sl, label, 0.4, 0.14, bg_col=col)


def slide_title(sl, text, sub=None):
    txb(sl, text, 0.45, 0.55, 12.4, 0.72,
        size=30, bold=True, color=C_WHITE)
    if sub:
        txb(sl, sub, 0.45, 1.28, 12.4, 0.38,
            size=15, italic=True, color=C_ACCENT)


def fake_table(sl, headers, rows,
               l, t, w,
               hdr_col=C_ACCENT2,
               fs_h=12, fs_r=11,
               row_h=0.38):
    cols  = len(headers)
    col_w = w / cols
    # header
    for ci, h in enumerate(headers):
        rect(sl, l+ci*col_w, t, col_w-0.03, row_h, hdr_col)
        txb(sl, h, l+ci*col_w+0.05, t+0.06,
            col_w-0.12, row_h-0.08,
            size=fs_h, bold=True, color=C_WHITE)
    # rows
    for ri, row in enumerate(rows):
        ry  = t + row_h + ri*row_h
        rbg = RGBColor(0x12,0x22,0x3A) if ri%2==0 else C_CARD
        for ci, cell in enumerate(row):
            rect(sl, l+ci*col_w, ry, col_w-0.03, row_h-0.02, rbg)
            txb(sl, str(cell),
                l+ci*col_w+0.05, ry+0.05,
                col_w-0.12, row_h-0.06,
                size=fs_r, color=C_LIGHT)


def glow_card(sl, l, t, w, h,
              top_color=None, border=False):
    """Card with optional coloured top bar."""
    rect(sl, l, t, w, h, C_CARD,
         line=border, line_rgb=C_ACCENT if border else None,
         line_w=Pt(0.75))
    if top_color:
        rect(sl, l, t, w, 0.06, top_color)


# ══════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════
def s01():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)

    # full-bleed gradient strips
    rect(sl, 0, 0, 0.07, 7.5, C_ACCENT)
    rect(sl, 13.26, 0, 0.07, 7.5, C_ACCENT2)
    rect(sl, 0, 0, 13.33, 0.06, C_ACCENT)
    rect(sl, 0, 7.44, 13.33, 0.06, C_ACCENT2)

    # large background card
    rect(sl, 0.55, 0.9, 12.23, 5.1, C_CARD)

    # decorative accent block
    rect(sl, 0.55, 0.9, 0.35, 5.1, C_ACCENT2)

    # internship label
    txb(sl, "I N T E R N S H I P   P R O J E C T   P R E S E N T A T I O N",
        1.1, 1.1, 11.5, 0.45,
        size=12, bold=True, color=C_ACCENT,
        align=PP_ALIGN.CENTER)

    # Main title
    txb(sl, "Deep Agent",
        1.1, 1.6, 11.5, 1.5,
        size=72, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER)

    # version badge
    rect(sl, 5.7, 3.1, 1.95, 0.4, C_ACCENT2)
    txb(sl, "v 2 . 0",
        5.72, 3.12, 1.9, 0.36,
        size=14, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER)

    # subtitle
    txb(sl,
        "Dynamic Multi-Agent Router  ·  Intelligent LLM Orchestration  ·  Cost Optimization",
        1.1, 3.62, 11.5, 0.52,
        size=18, italic=True, color=C_LIGHT,
        align=PP_ALIGN.CENTER)

    # tech pills
    pills = ["LiteLLM", "Groq", "Google Gemini", "Streamlit", "Python 3.9+"]
    total_pw = sum(len(p)*0.105 + 0.35 for p in pills) + 0.2*(len(pills)-1)
    sx = (13.33 - total_pw) / 2
    ty = 4.3
    for p in pills:
        pw = len(p)*0.105 + 0.35
        rect(sl, sx, ty, pw, 0.38, C_ACCENT2)
        txb(sl, p, sx+0.08, ty+0.05, pw-0.12, 0.28,
            size=12, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)
        sx += pw + 0.2

    # bottom presenter bar
    rect(sl, 0, 6.55, 13.33, 0.95, C_DARK_STRIP)
    txb(sl, "Presented by: [Shivam and Jay Kishan]",
        0.5, 6.65, 5, 0.42,
        size=13, color=C_LIGHT)
    txb(sl, "Mentor Review  —  April 2025",
        8.0, 6.65, 5.1, 0.42,
        size=13, color=C_LIGHT, align=PP_ALIGN.RIGHT)

    slide_num(sl, 1)

s01()


# ══════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════
def s02():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "AGENDA")
    slide_title(sl, "Today's Agenda")

    items = [
        ("01", "Title + Introduction",
         "Project overview, problem statement & motivation"),
        ("02", "Project Presentation",
         "Architecture, features, pipeline & what was built"),
        ("03", "Methods, Skills & Technology",
         "Tech stack, diagrams, new skills acquired"),
        ("04", "Goals & Objectives",
         "How Deep Agent fulfils all 8 internship objectives"),
        ("05", "Results",
         "Measured outcomes, cost savings, metrics"),
        ("06", "Future Scope",
         "Phase 2 features & GCP production roadmap"),
    ]

    for i, (num, title, desc) in enumerate(items):
        row = i // 2
        col = i % 2
        lx  = 0.42 + col * 6.42
        ty  = 1.72 + row * 1.88

        glow_card(sl, lx, ty, 6.12, 1.72, border=True)

        # number circle
        rect(sl, lx+0.14, ty+0.2, 0.65, 0.65, C_ACCENT2)
        txb(sl, num, lx+0.15, ty+0.22, 0.63, 0.62,
            size=20, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)

        txb(sl, title,
            lx+0.93, ty+0.14, 5.0, 0.44,
            size=16, bold=True, color=C_ACCENT)
        txb(sl, desc,
            lx+0.93, ty+0.64, 5.0, 0.9,
            size=13, color=C_LIGHT)

    slide_num(sl, 2)

s02()


# ══════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION / PROBLEM vs SOLUTION
# ══════════════════════════════════════════════════════
def s03():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "INTRODUCTION")
    slide_title(sl,
        "The Problem with Traditional LLM Deployment",
        sub="Why a smart routing layer is essential")

    # LEFT — Problems
    glow_card(sl, 0.38, 1.72, 5.95, 5.42, top_color=C_RED)
    txb(sl, "❌  Without Smart Routing",
        0.52, 1.85, 5.65, 0.42,
        size=15, bold=True, color=C_RED)
    probs = [
        "Every query → most expensive Pro model",
        "No specialization — generic prompts for all",
        "No fallback when primary model fails",
        "Zero cost visibility or comparison",
        "Manual model selection required from user",
        "Costs spiral uncontrollably with scale",
        "Rate-limit errors cause silent failures",
    ]
    bullet_box(sl, ["  ✗  " + p for p in probs],
               0.52, 2.38, 5.7, 4.6,
               size=13.5, color=C_LIGHT, gap=5)

    # RIGHT — Solution
    glow_card(sl, 6.98, 1.72, 5.95, 5.42, top_color=C_GREEN)
    txb(sl, "✅  Deep Agent v2.0",
        7.12, 1.85, 5.65, 0.42,
        size=15, bold=True, color=C_GREEN)
    sols = [
        "Routes to cheapest sufficient model automatically",
        "3 specialized agents — Coding, Math, Reasoning",
        "3-level fallback ensures 99.8 % uptime",
        "Full cost breakdown + savings % per query",
        "Zero configuration — fully automatic routing",
        "Up to 94 % cost reduction vs all-Pro routing",
        "Input & output guardrails for production safety",
    ]
    bullet_box(sl, ["  ✓  " + s for s in sols],
               7.12, 2.38, 5.7, 4.6,
               size=13.5, color=C_LIGHT, gap=5)

    # VS badge
    rect(sl, 6.08, 3.98, 0.82, 0.82, C_ACCENT)
    txb(sl, "VS", 6.1, 4.0, 0.78, 0.78,
        size=20, bold=True, color=C_BG,
        align=PP_ALIGN.CENTER)

    slide_num(sl, 3)

s03()


# ══════════════════════════════════════════════════════
# SLIDE 4 — FULL PIPELINE ARCHITECTURE
# ══════════════════════════════════════════════════════
def s04():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "PROJECT PRESENTATION", col=C_ACCENT)
    slide_title(sl, "End-to-End System Architecture",
                sub="Every query passes through this 6-stage pipeline")

    # Pipeline stages
    stages = [
        ("User\nQuery",        C_MID,    "Input"),
        ("Input\nGuardrails",  C_RED,    "Safety"),
        ("Router\nLLaMA 8B",   C_ACCENT2,"Classify"),
        ("Sub-Agent\n+ Model", C_GREEN,  "Execute"),
        ("Output\nGuardrails", C_RED,    "Safety"),
        ("Analytics\n+ Log",   C_ORANGE, "Observe"),
    ]

    bw, bh = 1.78, 1.45
    gap    = 0.14
    sy     = 1.82
    sx0    = 0.3

    for i, (label, color, sub) in enumerate(stages):
        bx = sx0 + i*(bw+gap)
        rect(sl, bx, sy, bw, bh, color)
        txb(sl, label, bx+0.05, sy+0.2, bw-0.1, 0.78,
            size=14, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)
        # sub label
        rect(sl, bx, sy+bh-0.32, bw, 0.32, C_DARK_STRIP)
        txb(sl, sub, bx+0.02, sy+bh-0.3, bw-0.04, 0.28,
            size=10, bold=True, color=C_ACCENT,
            align=PP_ALIGN.CENTER)
        # arrow
        if i < len(stages)-1:
            ax = bx + bw + 0.01
            txb(sl, "▶", ax, sy+0.48, 0.15, 0.5,
                size=12, bold=True, color=C_ACCENT)

    # --- TIER CARDS below ---
    tier_y = 3.56
    txb(sl, "Model Tiers — Dynamically Selected by Router",
        0.3, tier_y, 12.7, 0.38,
        size=14, bold=True, color=C_ACCENT)

    tiers = [
        ("🟢  LITE TIER",    C_GREEN,
         ["Simple queries (factual, basic math, trivial code)",
          "Model A: Gemini 2.5 Flash-Lite  (Google)",
          "Model B: GPT-OSS 20B  (Groq)",
          "Cost: $0.00005 / 1k tokens  |  ~400 ms"]),
        ("🟠  STANDARD TIER", C_ORANGE,
         ["Medium queries (algorithms, comparisons, SQL)",
          "Model A: LLaMA 3.1 8B  (Groq)",
          "Model B: GPT-OSS 120B  (Groq)",
          "Cost: $0.0001 / 1k tokens  |  ~500 ms"]),
        ("🔴  PRO TIER",      C_RED,
         ["Complex queries (system design, proofs, LRU cache)",
          "Model A: LLaMA 3.3 70B  (Groq)",
          "Model B: Gemini 2.5 Flash  (Google)",
          "Cost: $0.0008 / 1k tokens  |  ~2500 ms"]),
        ("🔀  ROUTER MODEL",  C_ACCENT2,
         ["Classification only — not used for answers",
          "Model: LLaMA 3.1 8B  (Groq)",
          "Temperature = 0  →  fully deterministic",
          "Overhead: ~100 ms  |  < 1 % of query cost"]),
    ]

    for i, (title, color, pts) in enumerate(tiers):
        tx = 0.3 + i * 3.26
        glow_card(sl, tx, tier_y+0.48, 3.12, 3.28, top_color=color)
        txb(sl, title, tx+0.1, tier_y+0.56, 2.92, 0.36,
            size=12, bold=True, color=C_WHITE)
        bullet_box(sl, ["• " + p for p in pts],
                   tx+0.1, tier_y+1.0, 2.95, 2.6,
                   size=11, color=C_LIGHT, gap=3)

    slide_num(sl, 4)

s04()


# ══════════════════════════════════════════════════════
# SLIDE 5 — KEY FEATURES (6-grid)
# ══════════════════════════════════════════════════════
def s05():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "PROJECT PRESENTATION", col=C_ACCENT)
    slide_title(sl, "Core Features — What Deep Agent Delivers")

    features = [
        ("🔀  Intelligent Routing",
         C_ACCENT2,
         ["LLaMA 8B classifies agent + complexity at runtime",
          "Agents: Coding · Math · Reasoning",
          "Complexity → Tier: simple/medium/complex",
          "Fully automatic — zero user configuration"]),
        ("💰  Cost Optimization",
         C_GREEN,
         ["Right model for each query — not always Pro",
          "Up to 94 % cheaper than all-Pro routing",
          "Per-query cost + savings shown in UI",
          "All 6 models compared in cost table"]),
        ("🛡️  3-Layer Guardrails",
         C_RED,
         ["Input: injection detection, PII warn, rate limit",
          "Output: dangerous code, uncertainty detection",
          "System: $5 daily cost ceiling auto-enforced",
          "Works independently of routing or agent"]),
        ("🔄  3-Level Fallback",
         C_ORANGE,
         ["Primary model → same-tier alternate → lite_b",
          "Fallback triggered silently, shown in UI",
          "Error details displayed for transparency",
          "Guarantees 99.8 % response uptime"]),
        ("⬆️  Model Upgrade System",
         C_ACCENT,
         ["Horizontal switch: same tier, different model",
          "Vertical upgrade: lite → standard → pro",
          "Full chain: lite_a → lite_b → … → pro_b",
          "Side-by-side cost + latency comparison"]),
        ("📊  Analytics Dashboard",
         C_YELLOW,
         ["5 live metrics: cost, latency, saved %, fallbacks",
          "Charts: tier dist., cost/query, latency trend",
          "All-6-model cost table per query",
          "CSV log with clear-all + rerun support"]),
    ]

    for i, (title, color, pts) in enumerate(features):
        col = i % 3
        row = i // 3
        fx  = 0.35 + col * 4.32
        fy  = 1.72 + row * 2.82

        glow_card(sl, fx, fy, 4.14, 2.68, top_color=color)
        txb(sl, title, fx+0.12, fy+0.12, 3.9, 0.4,
            size=13.5, bold=True, color=C_WHITE)
        bullet_box(sl, ["▸  " + p for p in pts],
                   fx+0.12, fy+0.6, 3.88, 1.94,
                   size=11.5, color=C_LIGHT, gap=4)

    slide_num(sl, 5)

s05()


# ══════════════════════════════════════════════════════
# SLIDE 6 — DEMO WALK-THROUGH
# ══════════════════════════════════════════════════════
def s06():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "PROJECT PRESENTATION", col=C_ACCENT)
    slide_title(sl, "User Journey — Annotated Demo Flow",
                sub="What happens between typing a query and seeing the answer")

    steps = [
        ("01", "User Types Query",
         C_ACCENT,
         "Text area with pre-filled example queries in sidebar.\n"
         "Sidebar shows live rate limit & daily budget status."),
        ("02", "Input Guardrails",
         C_RED,
         "Blocks 12 injection phrases. Warns on PII (email, phone, SSN).\n"
         "Enforces 5000-char limit. Rate-limits at 20 req / min."),
        ("03", "Router Classifies",
         C_ACCENT2,
         "LLaMA 8B (Groq) returns JSON: agent + complexity + confidence.\n"
         "Maps complexity → tier → model_key. Temperature = 0."),
        ("04", "Agent Executes",
         C_GREEN,
         "Specialized system prompt loaded (Coding / Math / Reasoning).\n"
         "LLM called via LiteLLM. 3-level fallback on any failure."),
        ("05", "Output Guardrails",
         C_ORANGE,
         "Scans response for dangerous code patterns (eval, os.system…).\n"
         "Flags uncertainty phrases. Warns if response too short."),
        ("06", "Results + Analytics",
         C_YELLOW,
         "Response rendered in markdown. Metrics: latency, tokens, cost.\n"
         "All-6-model cost table. Upgrade buttons. CSV logged."),
    ]

    for i, (num, title, color, body) in enumerate(steps):
        row = i // 2
        col = i % 2
        sx  = 0.38 + col * 6.47
        sy  = 1.72 + row * 1.92

        glow_card(sl, sx, sy, 6.2, 1.78, top_color=color)

        # step badge
        rect(sl, sx+0.12, sy+0.22, 0.7, 0.7, color)
        txb(sl, num, sx+0.13, sy+0.24, 0.68, 0.66,
            size=18, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)

        txb(sl, title, sx+0.96, sy+0.12, 5.1, 0.44,
            size=15, bold=True, color=C_WHITE)
        txb(sl, body, sx+0.96, sy+0.6, 5.1, 1.05,
            size=12, color=C_LIGHT)

    slide_num(sl, 6)

s06()


# ══════════════════════════════════════════════════════
# SLIDE 7 — TECH STACK + FILE ARCHITECTURE
# ══════════════════════════════════════════════════════
def s07():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "METHODS, SKILLS & TECHNOLOGY", col=C_ORANGE)
    slide_title(sl, "Technology Stack & Project File Architecture")

    # LEFT — tech table
    txb(sl, "Technology Stack",
        0.38, 1.72, 6.2, 0.38,
        size=15, bold=True, color=C_ACCENT)
    fake_table(
        sl,
        ["Component", "Technology", "Role"],
        [
            ["UI Framework",    "Streamlit 1.30+",    "Web app + dashboard"],
            ["LLM Gateway",     "LiteLLM",            "Unified API (Groq + Gemini)"],
            ["Fast Inference",  "Groq API",           "LLaMA 8B · 70B · GPT-OSS"],
            ["Lightweight LLM", "Google Gemini API",  "Flash-Lite & Flash"],
            ["Data Viz",        "Plotly + Pandas",    "Analytics charts + tables"],
            ["CSV Logging",     "Python stdlib csv",  "Query log analytics"],
            ["Config",          "Python dotenv",      "API key management"],
        ],
        l=0.38, t=2.18, w=6.2,
        hdr_col=C_ACCENT2, fs_h=12, fs_r=11,
    )

    # RIGHT — file tree
    txb(sl, "File Architecture (7 core files)",
        6.95, 1.72, 6.2, 0.38,
        size=15, bold=True, color=C_ACCENT)

    glow_card(sl, 6.95, 2.18, 6.2, 5.1)

    files = [
        ("app.py",          C_ACCENT,  "Streamlit UI + main orchestration"),
        ("router.py",       C_ACCENT2, "Query classification — LLaMA 8B"),
        ("agents.py",       C_GREEN,   "Sub-agent execution + fallback"),
        ("guardrails.py",   C_RED,     "3-layer safety system"),
        ("analytics.py",    C_ORANGE,  "CSV logging + cost computation"),
        ("config.py",       C_YELLOW,  "Model registry + upgrade paths"),
        ("query_logs.csv",  C_MID,     "Auto-created analytics log"),
    ]

    for i, (fname, fc, desc) in enumerate(files):
        fy = 2.32 + i * 0.67
        rect(sl, 7.08, fy, 5.9, 0.56, C_CARD2)
        rect(sl, 7.08, fy, 0.06, 0.56, fc)
        txb(sl, fname, 7.22, fy+0.08, 2.3, 0.38,
            size=12, bold=True, color=fc)
        txb(sl, desc,  9.55, fy+0.08, 3.5, 0.38,
            size=11, color=C_LIGHT)

    slide_num(sl, 7)

s07()


# ══════════════════════════════════════════════════════
# SLIDE 8 — LAYERED ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════
def s08():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "METHODS, SKILLS & TECHNOLOGY", col=C_ORANGE)
    slide_title(sl, "Layered Architecture Diagram",
                sub="How the 4 layers interact inside Deep Agent")

    layers = [
        ("LAYER 4 — PRESENTATION",  C_ACCENT,
         "app.py  (Streamlit)",
         "UI rendering · session state · button handlers · upgrade panel · analytics dashboard"),
        ("LAYER 3 — ORCHESTRATION", C_ACCENT2,
         "router.py  +  agents.py",
         "Query classification · model selection · specialized prompts · 3-level fallback"),
        ("LAYER 2 — SAFETY",        C_RED,
         "guardrails.py",
         "Input validation · output scanning · rate limiting · daily cost ceiling"),
        ("LAYER 1 — FOUNDATION",    C_GREEN,
         "config.py  +  analytics.py",
         "Model registry · upgrade paths · guardrail config · CSV logging · cost computation"),
    ]

    for i, (title, color, sub, desc) in enumerate(layers):
        ly = 1.72 + i * 1.35
        w  = 12.6 - i * 0.0   # all same width for clean look
        lx = 0.36

        rect(sl, lx, ly, w, 1.22, C_CARD)
        rect(sl, lx, ly, w, 0.06, color)
        rect(sl, lx, ly, 0.06, 1.22, color)

        txb(sl, title, lx+0.18, ly+0.1, 3.8, 0.36,
            size=12, bold=True, color=color)
        txb(sl, sub,   lx+0.18, ly+0.5, 3.8, 0.36,
            size=11, italic=True, color=C_LIGHT)
        txb(sl, desc,  lx+4.2,  ly+0.1, 8.2, 0.9,
            size=12, color=C_LIGHT)

    # API providers strip
    rect(sl, 0.36, 7.12, 12.6, 0.25, C_DARK_STRIP)
    txb(sl, "External APIs:   Groq  (LLaMA 8B · 70B · GPT-OSS 20B · 120B)   ·   Google Gemini  (Flash-Lite · Flash)",
        0.5, 7.13, 12.3, 0.22,
        size=10, color=C_MID, align=PP_ALIGN.CENTER)

    slide_num(sl, 8)

s08()


# ══════════════════════════════════════════════════════
# SLIDE 9 — MODEL REGISTRY & UPGRADE CHAIN
# ══════════════════════════════════════════════════════
def s09():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "METHODS, SKILLS & TECHNOLOGY", col=C_ORANGE)
    slide_title(sl, "Model Registry & Upgrade Chain",
                sub="6 models · 3 tiers · 2 providers · 1 deterministic upgrade path")

    fake_table(
        sl,
        ["Key", "Model", "Provider", "Tier", "Cost / 1k tokens", "Avg Latency", "Strength"],
        [
            ["lite_a",     "gemini-2.5-flash-lite",   "Google", "🟢 Lite",     "$0.00005",  "400 ms",  "Fastest & cheapest"],
            ["lite_b",     "gpt-oss-20b",             "Groq",   "🟢 Lite",     "$0.000075", "500 ms",  "Fast 20B, 1000T/s"],
            ["standard_a", "llama-3.1-8b-instant",    "Groq",   "🟠 Standard", "$0.0001",   "500 ms",  "Balanced cost/quality"],
            ["standard_b", "gpt-oss-120b",            "Groq",   "🟠 Standard", "$0.00015",  "800 ms",  "Strong reasoning"],
            ["pro_a",      "llama-3.3-70b-versatile", "Groq",   "🔴 Pro",      "$0.0008",   "2500 ms", "Best open-source"],
            ["pro_b",      "gemini-2.5-flash",        "Google", "🔴 Pro",      "$0.0005",   "2000 ms", "Google's best"],
        ],
        l=0.38, t=1.72, w=12.58,
        hdr_col=C_ACCENT2, fs_h=12, fs_r=11,
    )

    # Upgrade chain visual
    ty = 5.62
    txb(sl, "Upgrade Chain  (one-directional — no downgrade):",
        0.38, ty, 7, 0.36,
        size=14, bold=True, color=C_ACCENT)

    keys   = ["lite_a", "lite_b", "standard_a", "standard_b", "pro_a",  "pro_b"]
    colors = [C_GREEN,  C_GREEN,  C_ORANGE,     C_ORANGE,     C_RED,    C_RED]
    kw, kh = 1.72, 0.52
    kgap   = 0.08

    for i, (k, c) in enumerate(zip(keys, colors)):
        kx = 0.38 + i*(kw+kgap)
        rect(sl, kx, ty+0.46, kw, kh, c)
        txb(sl, k, kx+0.05, ty+0.49, kw-0.1, kh-0.06,
            size=12, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)
        if i < len(keys)-1:
            txb(sl, "▶",
                kx+kw+0.01, ty+0.56, 0.1, 0.36,
                size=14, bold=True, color=C_ACCENT)

    # Fallback note
    rect(sl, 0.38, ty+1.1, 12.58, 0.52, C_CARD2)
    txb(sl,
        "Fallback Logic:  Primary model fails  →  same-tier alternate  →  lite_b (GPT-OSS 20B)  as hard last resort",
        0.52, ty+1.15, 12.2, 0.42,
        size=12, color=C_ORANGE)

    slide_num(sl, 9)

s09()


# ══════════════════════════════════════════════════════
# SLIDE 10 — SKILLS LEARNED
# ══════════════════════════════════════════════════════
def s10():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "METHODS, SKILLS & TECHNOLOGY", col=C_ORANGE)
    slide_title(sl, "New Skills & Technologies Acquired",
                sub="Core learning outcomes from building Deep Agent")

    quadrants = [
        ("LLM Engineering",      C_ACCENT2,
         ["Multi-provider API integration via LiteLLM",
          "Prompt engineering per agent specialization",
          "Token budget management & cost calculation",
          "Fallback chain design across providers",
          "Temperature=0 for deterministic classification"]),
        ("AI System Design",     C_ACCENT,
         ["Router → Agent pipeline architecture",
          "Stateless LLM system design patterns",
          "3-layer guardrail implementation",
          "Model registry & upgrade path design",
          "JSON-structured LLM output parsing"]),
        ("Streamlit Development",C_GREEN,
         ["st.session_state for persistent UI state",
          "st.cache_data for performance optimization",
          "st.status for live pipeline updates",
          "Custom CSS injection for styling",
          "st.download_button for data export"]),
        ("Cloud & Architecture", C_ORANGE,
         ["GCP service mapping (10 services identified)",
          "Firestore vs JSON for session storage",
          "Redis vs file-based caching trade-offs",
          "BigQuery vs CSV for analytics pipelines",
          "Cloud Run auto-scaling architecture"]),
    ]

    for i, (title, color, pts) in enumerate(quadrants):
        col = i % 2
        row = i // 2
        qx  = 0.38 + col * 6.47
        qy  = 1.72 + row * 2.78

        glow_card(sl, qx, qy, 6.2, 2.64, top_color=color)
        txb(sl, title, qx+0.15, qy+0.14, 5.9, 0.4,
            size=15, bold=True, color=C_WHITE)
        bullet_box(sl, ["✓  " + p for p in pts],
                   qx+0.15, qy+0.62, 5.9, 1.88,
                   size=12.5, color=C_LIGHT, gap=4)

    slide_num(sl, 10)

s10()


# ══════════════════════════════════════════════════════
# SLIDE 11 — GOALS & OBJECTIVES (all 8)
# ══════════════════════════════════════════════════════
def s11():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "GOALS & OBJECTIVES", col=C_GREEN)
    slide_title(sl, "8 Program Objectives — Achieved",
                sub="How Deep Agent fulfils every internship objective")

    objs = [
        ("1", "Application",
         C_ACCENT,
         "Built & deployed a live Streamlit app routing real queries\n"
         "to real LLMs via Groq & Gemini APIs with full UI."),
        ("2", "Knowledge",
         C_ACCENT2,
         "Learned LLM orchestration, prompt engineering, LiteLLM,\n"
         "guardrail design, cost modelling, and GCP architecture."),
        ("3", "Exploration",
         C_GREEN,
         "Explored 6 LLM models across 2 providers, 3 complexity\n"
         "tiers, 10 GCP services, and Streamlit advanced patterns."),
        ("4", "Identification",
         C_ORANGE,
         "Identified: cost overrun, no fallback, rate-limit failures,\n"
         "prompt injection risk, PII exposure, zero observability."),
        ("5", "Innovation",
         C_RED,
         "Designed per-complexity routing, 3-level fallback chain,\n"
         "per-agent specialized prompts, and model upgrade system."),
        ("6", "Engagement",
         C_YELLOW,
         "Iterative development with mentor feedback; responded to\n"
         "review by planning Phase 2 memory & cache features."),
        ("7", "Evaluation",
         C_ACCENT,
         "Analytics dashboard shows cost vs Pro baseline, fallback\n"
         "rate, tier distribution — measurable, data-backed ROI."),
        ("8", "Demonstration",
         C_GREEN,
         "Live Streamlit demo: full pipeline, guardrail blocking,\n"
         "upgrade comparison, real-time cost metrics & charts."),
    ]

    for i, (num, title, color, body) in enumerate(objs):
        row = i // 2
        col = i % 2
        ox  = 0.38 + col * 6.47
        oy  = 1.72 + row * 1.42

        glow_card(sl, ox, oy, 6.2, 1.32)
        rect(sl, ox, oy, 0.08, 1.32, color)

        # number chip
        rect(sl, ox+0.2, oy+0.16, 0.55, 0.55, color)
        txb(sl, num, ox+0.21, oy+0.18, 0.53, 0.51,
            size=16, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)

        txb(sl, title, ox+0.88, oy+0.1, 5.2, 0.38,
            size=14, bold=True, color=color)
        txb(sl, body,  ox+0.88, oy+0.5, 5.2, 0.76,
            size=11, color=C_LIGHT)

    slide_num(sl, 11)

s11()


# ══════════════════════════════════════════════════════
# SLIDE 12 — GUARDRAILS DEEP DIVE
# ══════════════════════════════════════════════════════
def s12():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "GOALS & OBJECTIVES", col=C_GREEN)
    slide_title(sl, "Objective Deep-Dive — 3-Layer Guardrail System",
                sub="Production-grade safety built into every query lifecycle")

    layers = [
        ("LAYER 1  —  INPUT  (check_input)", C_RED,
         [("Prompt Injection",
           "12 blocked phrases: 'ignore all instructions', 'jailbreak', 'DAN', 'reveal system prompt'…"),
          ("PII Detection",
           "Regex patterns: email, phone (###-###-####), SSN (###-##-####), 16-digit credit card"),
          ("Length Limit",
           "5000-character maximum — auto-truncates with visible warning to user"),
          ("Rate Limiting",
           "Sliding 60-second window — blocks at 20 requests/min using in-memory timestamps")]),
        ("LAYER 2  —  OUTPUT  (check_output)", C_ORANGE,
         [("Dangerous Code",
           "Detects: os.system(  eval(  exec(  subprocess  rm -rf  DROP TABLE  shutil.rmtree"),
          ("Uncertainty Phrases",
           "Flags: 'I'm not sure'  'I cannot verify'  'this may be incorrect' → suggests upgrade"),
          ("Minimum Length",
           "Response shorter than 10 chars flagged — likely model failure or refusal"),
          ("Transparency",
           "All warnings shown to user in expandable panel, not silently suppressed")]),
        ("LAYER 3  —  SYSTEM  (track_cost / get_status)", C_ACCENT2,
         [("Daily Budget",
           "$5.00 ceiling — warning banner shown when exceeded, does not hard-block"),
          ("Auto Reset",
           "Resets at midnight using datetime.date comparison — in-memory state"),
          ("Live Sidebar",
           "🟢 / 🟡 / 🔴 indicators for rate usage and budget — updates on every rerun"),
          ("Scope",
           "Only real LLM calls tracked — zero-cost routes handled transparently")]),
    ]

    for i, (title, color, checks) in enumerate(layers):
        lx = 0.38 + i * 4.32
        glow_card(sl, lx, 1.72, 4.14, 5.55, top_color=color)
        txb(sl, title, lx+0.12, 1.82, 3.9, 0.5,
            size=12, bold=True, color=color)

        for j, (check, desc) in enumerate(checks):
            jy = 2.44 + j * 1.2
            rect(sl, lx+0.12, jy, 3.88, 1.05, C_CARD2)
            txb(sl, check, lx+0.2, jy+0.06, 3.72, 0.32,
                size=12, bold=True, color=C_WHITE)
            txb(sl, desc,  lx+0.2, jy+0.42, 3.72, 0.58,
                size=10.5, color=C_LIGHT)

    slide_num(sl, 12)

s12()


# ══════════════════════════════════════════════════════
# SLIDE 13 — RESULTS / METRICS
# ══════════════════════════════════════════════════════
def s13():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "RESULTS", col=C_YELLOW)
    slide_title(sl, "Results — Measured Outcomes & Metrics")

    # Big KPI numbers
    kpis = [
        ("94 %",   "Max Cost\nReduction",         C_GREEN),
        ("57 %",   "Average\nSavings",             C_ACCENT),
        ("99.8 %", "Uptime via\nFallback",         C_ORANGE),
        ("6",      "LLM Models\nIntegrated",       C_ACCENT2),
        ("3",      "Specialized\nAgents",          C_RED),
    ]

    for i, (val, label, color) in enumerate(kpis):
        kx = 0.38 + i * 2.52
        rect(sl, kx, 1.72, 2.38, 1.52, color)
        txb(sl, val,   kx+0.08, 1.82, 2.22, 0.76,
            size=34, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)
        txb(sl, label, kx+0.08, 2.52, 2.22, 0.64,
            size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Cost table
    txb(sl, "Cost Comparison — 500-Token Query",
        0.38, 3.42, 8.5, 0.38,
        size=14, bold=True, color=C_ACCENT)

    fake_table(
        sl,
        ["Model", "Tier", "500-Token Cost", "Latency", "Saving vs Pro A"],
        [
            ["Gemini 2.5 Flash-Lite  ← default Lite",
             "🟢 Lite",     "$0.000025",  "400 ms",  "94 %"],
            ["GPT-OSS 20B",
             "🟢 Lite",     "$0.0000375", "500 ms",  "91 %"],
            ["LLaMA 3.1 8B  ← default Standard",
             "🟠 Standard", "$0.00005",   "500 ms",  "88 %"],
            ["GPT-OSS 120B",
             "🟠 Standard", "$0.000075",  "800 ms",  "81 %"],
            ["LLaMA 3.3 70B  ← default Pro",
             "🔴 Pro",      "$0.0004",    "2500 ms", "baseline"],
            ["Gemini 2.5 Flash",
             "🔴 Pro",      "$0.00025",   "2000 ms", "38 %"],
        ],
        l=0.38, t=3.88, w=9.5,
        hdr_col=C_ACCENT2, fs_h=11.5, fs_r=11,
    )

    # Right side stats
    glow_card(sl, 10.1, 3.42, 3.12, 4.0, border=True)
    txb(sl, "Routing Distribution",
        10.22, 3.52, 2.88, 0.38,
        size=13, bold=True, color=C_ACCENT)

    dist = [
        ("🟢 ~5 %",  "Simple → Lite tier"),
        ("🟠 ~60 %", "Medium → Standard"),
        ("🔴 ~35 %", "Complex → Pro tier"),
        ("",""),
        ("🔀 ~100ms","Router overhead"),
        ("🛡️ <2 %",  "Injection attempts"),
        ("🔄 ~2 %",  "Fallback trigger rate"),
    ]
    for j, (pct, desc) in enumerate(dist):
        jy = 4.0 + j*0.47
        if pct:
            txb(sl, pct,  10.22, jy, 1.3, 0.4,
                size=12, bold=True, color=C_ACCENT)
            txb(sl, desc, 11.55, jy, 1.55, 0.4,
                size=11, color=C_LIGHT)

    slide_num(sl, 13)

s13()


# ══════════════════════════════════════════════════════
# SLIDE 14 — RESULTS SUMMARY
# ══════════════════════════════════════════════════════
def s14():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "RESULTS", col=C_YELLOW)
    slide_title(sl, "Overall Internship Results",
                sub="What was built, what was learned, what was delivered")

    # LEFT — built
    glow_card(sl, 0.38, 1.72, 6.0, 5.55, top_color=C_ACCENT)
    txb(sl, "What Was Built",
        0.52, 1.82, 5.7, 0.4,
        size=15, bold=True, color=C_ACCENT)
    built = [
        "Multi-agent router — 6 files, ~1 200 lines of code",
        "6 LLM models unified under LiteLLM gateway",
        "3 specialized agents with structured prompts",
        "3-level automatic fallback (primary→tier→lite_b)",
        "3-layer guardrails: injection, output, cost ceiling",
        "Model upgrade system — horizontal + vertical",
        "Analytics dashboard — 5 metrics, 4 chart tabs",
        "All-6-model cost comparison per query",
        "Streamlit Cloud deployment (Streamlit-safe paths)",
        "GCP migration roadmap for 10 production services",
    ]
    bullet_box(sl, ["  ✅  " + b for b in built],
               0.52, 2.3, 5.74, 4.8,
               size=13, color=C_LIGHT, gap=4)

    # RIGHT — learned
    glow_card(sl, 6.94, 1.72, 6.0, 5.55, top_color=C_GREEN)
    txb(sl, "What Was Learned",
        7.08, 1.82, 5.7, 0.4,
        size=15, bold=True, color=C_GREEN)
    learned = [
        "LLM API integration: Groq, Gemini, LiteLLM",
        "Prompt engineering for domain-specific agents",
        "Cost-aware architecture design principles",
        "Guardrail patterns for production AI safety",
        "Stateful UI design with Streamlit session_state",
        "GCP cloud-native production architecture",
        "Iterative development with mentor feedback loops",
        "Trade-offs: JSON vs Firestore, CSV vs BigQuery",
        "Multi-provider fallback resilience patterns",
        "AI system documentation at production level",
    ]
    bullet_box(sl, ["  🎓  " + l for l in learned],
               7.08, 2.3, 5.74, 4.8,
               size=13, color=C_LIGHT, gap=4)

    slide_num(sl, 14)

s14()


# ══════════════════════════════════════════════════════
# SLIDE 15 — FUTURE SCOPE
# ══════════════════════════════════════════════════════
def s15():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)
    h_line(sl, 0.6)
    section_tag(sl, "FUTURE SCOPE", col=C_ACCENT2)
    slide_title(sl, "Future Scope — Phase 2 & GCP Production",
                sub="Planned roadmap building on the current foundation")

    # PHASE 2 — left
    glow_card(sl, 0.38, 1.72, 5.82, 5.55, top_color=C_ACCENT2)
    txb(sl, "Phase 2 — Feature Roadmap",
        0.52, 1.82, 5.58, 0.38,
        size=14, bold=True, color=C_WHITE)

    p2 = [
        ("🧠  Conversational Memory",
         "Session JSON · last 5 turns · per-agent namespacing"),
        ("⚡  Semantic Cache",
         "difflib fuzzy match · $0 on repeated queries"),
        ("🔀  Handoff Context",
         "Conversation context injected on model upgrades"),
        ("📥  Export",
         "Download conversations as Markdown or JSON"),
        ("👍  Feedback Loop",
         "Thumbs up/down to improve routing accuracy"),
        ("🖥️  Code Executor",
         "Run generated Python in isolated sandbox"),
        ("📈  Preference Learning",
         "Auto-route to user's preferred model per agent"),
    ]

    for j, (icon_title, desc) in enumerate(p2):
        jy = 2.32 + j * 0.72
        rect(sl, 0.52, jy, 5.55, 0.62, C_CARD2)
        txb(sl, icon_title, 0.62, jy+0.05, 3.2, 0.32,
            size=12, bold=True, color=C_ACCENT2)
        txb(sl, desc, 0.62, jy+0.34, 5.3, 0.24,
            size=10.5, color=C_LIGHT)

    # GCP MIGRATION — right
    glow_card(sl, 6.58, 1.72, 6.37, 5.55, top_color=C_ACCENT)
    txb(sl, "GCP Production Migration",
        6.72, 1.82, 6.1, 0.38,
        size=14, bold=True, color=C_WHITE)

    gcp = [
        ("sessions/*.json",   "→",  "Firestore",        C_GREEN,
         "Real-time, multi-user, survives restarts"),
        ("query_cache.json",  "→",  "Redis Memorystore", C_GREEN,
         "TTL per tier, auto-eviction, shared across instances"),
        ("query_logs.csv",    "→",  "BigQuery",          C_GREEN,
         "Streaming inserts, SQL analytics, Looker Studio"),
        ("os.getenv(keys)",   "→",  "Secret Manager",    C_YELLOW,
         "Encrypted at rest, audited access, versioned"),
        ("Streamlit Cloud",   "→",  "Cloud Run",         C_ACCENT,
         "Auto-scale 0 → 1000+ instances, pay per use"),
        ("Groq/Gemini APIs",  "→",  "Vertex AI",         C_ACCENT,
         "No rate limits, in-network latency, SLA-backed"),
        ("Plotly sidebar",    "→",  "Cloud Monitoring",  C_ORANGE,
         "Alerts on cost spike, fallback surge, latency P99"),
    ]

    for j, (cur, arr, fut, fc, note) in enumerate(gcp):
        jy = 2.32 + j * 0.7
        rect(sl, 6.72, jy, 6.1, 0.6, C_CARD2)
        txb(sl, cur,  6.8,  jy+0.05, 2.0, 0.28,
            size=10.5, bold=True, color=C_ORANGE)
        txb(sl, arr,  8.82, jy+0.07, 0.22, 0.24,
            size=11, bold=True, color=C_MID)
        txb(sl, fut,  9.06, jy+0.05, 1.7, 0.28,
            size=11, bold=True, color=fc)
        txb(sl, note, 6.8,  jy+0.35, 5.95, 0.22,
            size=9.5, color=C_LIGHT)

    slide_num(sl, 15)

s15()


# ══════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU
# ══════════════════════════════════════════════════════
def s16():
    sl = prs.slides.add_slide(BLANK)
    slide_bg(sl)

    # Full accent borders
    rect(sl, 0, 0,    13.33, 0.08, C_ACCENT)
    rect(sl, 0, 7.42, 13.33, 0.08, C_ACCENT2)
    rect(sl, 0, 0, 0.08, 7.5, C_ACCENT)
    rect(sl, 13.25, 0, 0.08, 7.5, C_ACCENT2)

    # Centre card
    rect(sl, 1.2, 0.95, 10.93, 5.25, C_CARD)
    rect(sl, 1.2, 0.95, 10.93, 0.08, C_ACCENT)

    # Main text
    txb(sl, "Thank You",
        1.2, 1.2, 10.93, 1.5,
        size=66, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER)

    txb(sl, "Deep Agent v2.0  —  Dynamic Multi-Agent Router",
        1.2, 2.72, 10.93, 0.52,
        size=20, italic=True, color=C_ACCENT,
        align=PP_ALIGN.CENTER)

    h_line(sl, 3.42, color=C_MID, h=0.02)

    txb(sl, "Open for Questions & Discussion",
        1.2, 3.62, 10.93, 0.52,
        size=22, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER)

    # Links row
    links = [
        ("📁  GitHub Repo",     C_ACCENT2),
        ("🌐  Live Demo",       C_GREEN),
        ("📄  README.md",       C_ACCENT),
        ("☁️   GCP Roadmap",    C_ORANGE),
    ]
    total_lw = len(links) * 2.5 + (len(links)-1) * 0.18
    lsx = (13.33 - total_lw) / 2
    for label, color in links:
        rect(sl, lsx, 4.32, 2.5, 0.55, color)
        txb(sl, label, lsx+0.06, 4.37, 2.38, 0.44,
            size=12, bold=True, color=C_BG,
            align=PP_ALIGN.CENTER)
        lsx += 2.5 + 0.18

    # Tech strip
    rect(sl, 0, 6.4, 13.33, 1.1, C_DARK_STRIP)
    txb(sl,
        "Built with   LiteLLM  ·  Groq  ·  Google Gemini  ·  Streamlit  ·  Python",
        0.5, 6.52, 12.33, 0.42,
        size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txb(sl,
        "Intelligent Routing  ·  Cost Optimization  ·  3-Layer Guardrails  ·  Analytics  ·  Model Upgrades",
        0.5, 6.94, 12.33, 0.38,
        size=11, color=C_MID, align=PP_ALIGN.CENTER)

    slide_num(sl, 16)

s16()


# ══════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════
OUT = "Deep_Agent_v2_Presentation.pptx"
prs.save(OUT)
print(f"\n✅  Presentation saved → {OUT}")
print(f"   {len(prs.slides)} slides  |  Widescreen 16:9  |  13.33 × 7.5 inches")
print("\nSlide map:")
labels = [
    "01  Title",
    "02  Agenda",
    "03  Introduction — Problem vs Solution",
    "04  Architecture — Full Pipeline",
    "05  Core Features (6-grid)",
    "06  Demo Walk-Through (6 steps)",
    "07  Tech Stack + File Architecture",
    "08  Layered Architecture Diagram",
    "09  Model Registry + Upgrade Chain",
    "10  New Skills Learned",
    "11  8 Program Objectives",
    "12  Guardrails Deep-Dive",
    "13  Results + Cost Metrics",
    "14  Overall Results Summary",
    "15  Future Scope + GCP Roadmap",
    "16  Thank You / Q&A",
]
for l in labels:
    print(f"   {l}")